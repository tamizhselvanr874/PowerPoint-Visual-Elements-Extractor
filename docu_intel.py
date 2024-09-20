import streamlit as st  
from pptx import Presentation  
import io  
import os  
from pptx.enum.shapes import MSO_SHAPE_TYPE  
from docx import Document  
from docx.shared import Inches  
import fitz  # PyMuPDF  
from PIL import Image  
from io import BytesIO  
import requests  
  
# URL of your Azure function endpoint  
AZURE_FUNCTION_URL = 'https://doc2pdf.azurewebsites.net/api/HttpTrigger1'  
  
  
def ppt_to_pdf(ppt_file, pdf_file_path):  
    """Convert PPT to PDF using Azure Function"""  
    mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'  
    headers = {  
        "Content-Type": "application/octet-stream",  
        "Content-Type-Actual": mime_type  
    }  
    response = requests.post(AZURE_FUNCTION_URL, data=ppt_file, headers=headers)  
    if response.status_code == 200:  
        with open(pdf_file_path, 'wb') as pdf_out:  
            pdf_out.write(response.content)  
        return True  
    else:  
        st.error(f"File conversion failed with status code: {response.status_code}")  
        st.error(f"Response: {response.text}")  
        return False  
  
  
def is_image_of_interest(shape):  
    """Check if a shape contains an image in formats of interest"""  
    try:  
        if hasattr(shape, "image"):  
            image_ext = os.path.splitext(shape.image.filename)[1].lower()  
            if image_ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"]:  
                return image_ext  
    except Exception:  
        pass  
    return None  
  
  
def detect_image_slides(ppt_file):  
    """Detect slides containing images in the desired formats"""  
    ppt = Presentation(io.BytesIO(ppt_file))  
    image_slides = {}  
    for i, slide in enumerate(ppt.slides):  
        for shape in slide.shapes:  
            image_format = is_image_of_interest(shape)  
            if image_format:  
                slide_number = i + 1  
                image_slides[slide_number] = image_format  
                break  
    return image_slides  
  
  
def identify_visual_elements(ppt_file):  
    """Identify slides with visual elements"""  
    presentation = Presentation(io.BytesIO(ppt_file))  
    visual_slides = []  
    for slide_number, slide in enumerate(presentation.slides, start=1):  
        has_visual_elements = False  
        for shape in slide.shapes:  
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART,  
                                    MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE}:  
                has_visual_elements = True  
                break  
        if has_visual_elements:  
            visual_slides.append(slide_number)  
    return visual_slides  
  
  
def combine_slide_numbers(image_slides, visual_slides):  
    """Combine slide numbers from image slides and visual element slides"""  
    combined_slides = set(image_slides.keys()).union(set(visual_slides))  
    return sorted(list(combined_slides))  
  
  
def capture_slide_images(pdf_file, slide_numbers):  
    """Capture images from identified slides in the PDF"""  
    doc = fitz.open(pdf_file)  
    images = []  
    for slide_number in slide_numbers:  
        page = doc[slide_number - 1]  
        pix = page.get_pixmap()  
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  
        buffer = BytesIO()  
        img.save(buffer, format="PNG")  
        images.append({"slide_number": slide_number, "image": buffer.getvalue()})  
    return images  
  
  
def generate_word_doc(slide_images):  
    """Generate Word document with slide images"""  
    doc = Document()  
    doc.add_heading('Slides with Visual Elements', level=1)  
    for item in slide_images:  
        slide_number = item["slide_number"]  
        image_data = item["image"]  
        doc.add_heading(f'Slide {slide_number}', level=2)  
        image_stream = BytesIO(image_data)  
        doc.add_picture(image_stream, width=Inches(6))  
    doc_file = BytesIO()  
    doc.save(doc_file)  
    doc_file.seek(0)  
    return doc_file  
  
  
def main():  
    st.title("PowerPoint Visual Elements Extractor")  
  
    # Upload PowerPoint file  
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=['pptx'])  
  
    if uploaded_file is not None:  
        ppt_file = uploaded_file.read()  
        ppt_file_path = "uploaded_ppt.pptx"  
        pdf_file_path = "converted.pdf"  
  
        # Save the uploaded file to disk  
        with open(ppt_file_path, 'wb') as f:  
            f.write(ppt_file)  
  
        # Identify image slides  
        image_slides = detect_image_slides(ppt_file)  
        st.write(f"Slides with images: {list(image_slides.keys())}")  
  
        # Identify slides with visual elements  
        visual_slides = identify_visual_elements(ppt_file)  
        st.write(f"Slides with visual elements: {visual_slides}")  
  
        # Combine slide numbers  
        combined_slides = combine_slide_numbers(image_slides, visual_slides)  
        st.write(f"Combined slides with visuals: {combined_slides}")  
  
        # Convert PPT to PDF using Azure Function  
        if ppt_to_pdf(ppt_file, pdf_file_path):  
            # Capture slides  
            slide_images = capture_slide_images(pdf_file_path, combined_slides)  
  
            # Generate and download Word document  
            doc_file = generate_word_doc(slide_images)  
            st.download_button(  
                "Download Word Document",  
                data=doc_file,  
                file_name="slides_with_visuals.docx",  
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"  
            )  
  
  
if __name__ == "__main__":  
    main()  
